"""
Build ZPT Sales Deck — 10 slides, ZPT brand (Navy/Gold/Off-white).
Run: python3 build_sales_deck.py
Output: ZPT_Sales_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand tokens ──
NAVY = RGBColor(0x0C, 0x0C, 0x28)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x88, 0x88, 0x99)
LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0xA8)
NAVY_LIGHT = RGBColor(0x1A, 0x1A, 0x3E)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)

FONT_HEADING = "Cormorant Garamond"
FONT_BODY = "Plus Jakarta Sans"
FONT_MONO = "Source Code Pro"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # blank layout


# ── Helpers ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             font_name=FONT_BODY, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  font_name=FONT_BODY, color=WHITE, alignment=PP_ALIGN.LEFT,
                  line_spacing=1.4, bullet_color=None):
    """lines: list of (text, bold) tuples or strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, bold = line
        else:
            text, bold = line, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_card(slide, left, top, width, height, title, body, fill=WHITE,
             title_color=NAVY, body_color=MUTED, border_color=None,
             title_size=16, body_size=13):
    """Add a rounded card with title + body text."""
    card = add_rect(slide, left, top, width, height, fill, border_color)
    # Title
    add_text(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.5),
             title, font_size=title_size, font_name=FONT_BODY, color=title_color, bold=True)
    # Body
    add_text(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
             body, font_size=body_size, font_name=FONT_BODY, color=body_color)
    return card


def add_gold_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "AI-powered sales operations.",
         font_size=44, font_name=FONT_HEADING, color=WHITE, bold=False,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(2.9), Inches(10), Inches(1),
         "Managed for you.",
         font_size=44, font_name=FONT_HEADING, color=GOLD, bold=False,
         alignment=PP_ALIGN.CENTER)

add_gold_divider(s, Inches(5.5), Inches(4.2), Inches(2.3))

add_text(s, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
         "zpteam.ai  |  request@zpteam.ai",
         font_size=14, font_name=FONT_BODY, color=MUTED,
         alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 2 — The Problem -> The Solution
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "The problem. The solution.",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

# Left card — The Problem
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
add_text(s, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
         "The reality for most sales teams",
         font_size=18, font_name=FONT_BODY, color=NAVY, bold=True)

problem_lines = [
    "+  Manual enrichment — copy-pasting between tabs for hours",
    "+  Scattered data across spreadsheets, CRMs, and inboxes",
    "+  Inconsistent prospect profiles — everyone does it differently",
    "+  Hours spent per prospect on research that goes stale",
    "+  No time left for actual selling",
]
add_multiline(s, Inches(1.1), Inches(2.5), Inches(5), Inches(3.5),
              problem_lines, font_size=14, color=MUTED, line_spacing=1.8)

# Right card — The Solution
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.8), NAVY)
add_text(s, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5),
         "What ZPT does",
         font_size=18, font_name=FONT_BODY, color=GOLD, bold=True)

add_text(s, Inches(7.3), Inches(2.5), Inches(4.8), Inches(2.5),
         "You share your sales context and data.\n\nZPT sources, enriches, qualifies, and delivers results back through your CRM.",
         font_size=16, font_name=FONT_BODY, color=WHITE, line_spacing=1.6)

add_rect(s, Inches(7.3), Inches(4.8), Inches(4.8), Inches(1.2), NAVY_LIGHT)
add_text(s, Inches(7.6), Inches(4.95), Inches(4.3), Inches(0.8),
         "This isn't a tool you learn.\nIt's a service you receive.",
         font_size=15, font_name=FONT_BODY, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 3 — The Full Sales Funnel
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "The full sales funnel",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

# Source labels at top
sources = ["LinkedIn", "Google Maps", "Web", "Job boards", "Conferences", "..."]
for i, src in enumerate(sources):
    x = Inches(1.0 + i * 1.9)
    add_rect(s, x, Inches(1.6), Inches(1.6), Inches(0.45), NAVY)
    add_text(s, x, Inches(1.63), Inches(1.6), Inches(0.4),
             src, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)

# Arrow down
add_text(s, Inches(5.8), Inches(2.15), Inches(1.5), Inches(0.4),
         "\u25BC", font_size=18, color=GOLD, alignment=PP_ALIGN.CENTER)

# 4 funnel stages as horizontal cards
stages = [
    ("1", "Find Prospects",
     "Source leads from LinkedIn, directories, maps, conferences, and competitor intel."),
    ("2", "Enrich & Qualify",
     "Build full profiles with decision-maker data, company intel, and ICP scoring."),
    ("3", "Personalized Outreach",
     "Draft emails using the company's voice, positioning, and prospect research."),
    ("4", "Manage Pipeline",
     "Track activity, follow-ups, and sync with the existing CRM."),
]

for i, (num, title, desc) in enumerate(stages):
    y = Inches(2.6 + i * 1.05)
    # Shrink width slightly per stage (funnel effect)
    inset = i * 0.35
    card_w = Inches(11.5 - inset * 2)
    card_x = Inches(0.9 + inset)

    bg_color = NAVY if i % 2 == 0 else NAVY_LIGHT
    add_rect(s, card_x, y, card_w, Inches(0.85), bg_color)

    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, card_x + Inches(0.25), y + Inches(0.15), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_text(s, card_x + Inches(0.25), y + Inches(0.2), Inches(0.55), Inches(0.45),
             num, font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text(s, card_x + Inches(1.0), y + Inches(0.1), Inches(3), Inches(0.5),
             title, font_size=17, color=WHITE, bold=True)

    # Description
    add_text(s, card_x + Inches(4.2), y + Inches(0.15), card_w - Inches(4.7), Inches(0.6),
             desc, font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD))

# Bottom callout
add_rect(s, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.6), CREAM)
add_text(s, Inches(3.5), Inches(6.52), Inches(6.3), Inches(0.55),
         "ZPT handles every stage. You approve the output.",
         font_size=15, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 4 — What the Output Looks Like
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "What the output looks like",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

# ── LEFT: Simple Template ──
add_text(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "Simple Template", font_size=16, color=NAVY, bold=True)

# Table
rows_data = [
    ["Company", "Contact", "Role", "Email", "Status"],
    ["Acme Corp", "Jane Smith", "CRO", "jane@acme.com", "Verified"],
    ["Beta Inc", "Tom Brown", "VP Sales", "tom@beta.io", "Enriched"],
    ["Gamma SA", "Lisa Muller", "CEO", "lisa@gamma.de", "Verified"],
    ["Delta Ltd", "Mark Roberts", "Head of Sales", "mark@delta.co.uk", "Pending"],
    ["Echo GmbH", "Anna Weber", "CMO", "anna@echo.de", "Verified"],
]

tbl = s.shapes.add_table(len(rows_data), 5,
                          Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.8)).table
tbl.columns[0].width = Inches(1.1)
tbl.columns[1].width = Inches(1.2)
tbl.columns[2].width = Inches(1.1)
tbl.columns[3].width = Inches(1.5)
tbl.columns[4].width = Inches(0.9)

for row_idx, row_data in enumerate(rows_data):
    for col_idx, val in enumerate(row_data):
        cell = tbl.cell(row_idx, col_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_BODY
            if row_idx == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            else:
                para.font.color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else CREAM

# ── RIGHT: Advanced Template ──
add_text(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
         "Advanced Template", font_size=16, color=NAVY, bold=True)

# Advanced card
add_rect(s, Inches(7.2), Inches(2.0), Inches(5.3), Inches(4.5), WHITE, border_color=GOLD)
add_text(s, Inches(7.5), Inches(2.15), Inches(4.7), Inches(0.4),
         "Acme Corp", font_size=20, color=NAVY, bold=True)
add_text(s, Inches(7.5), Inches(2.6), Inches(4.7), Inches(0.35),
         "SaaS, 50-200 employees, Munich, Germany", font_size=12, color=MUTED)

advanced_fields = [
    ("ICP-Fit", "5 / 5"),
    ("Decision maker", "Jane Smith, CRO"),
    ("Pain signals", "Hiring SDR, manual CRM updates, low outreach volume"),
    ("Personalized hook", '"Saw your SDR job post on LinkedIn. ZPT could cover that function at a fraction of the cost..."'),
    ("Recommended channel", "LinkedIn + email sequence"),
]

y_pos = Inches(3.1)
for label, value in advanced_fields:
    add_text(s, Inches(7.5), y_pos, Inches(2), Inches(0.3),
             label, font_size=11, color=GOLD, bold=True)
    add_text(s, Inches(9.3), y_pos, Inches(3), Inches(0.55),
             value, font_size=11, color=NAVY)
    y_pos += Inches(0.55)

add_text(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5),
         "Templates are fully customizable. Each client defines their own fields, scoring criteria, and output format.",
         font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 5 — How It Works Under the Hood
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "How it works",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "One system. Reads your context, searches the web, updates your CRM, delivers structured output.",
         font_size=16, color=MUTED)

# Center: The folder
add_rect(s, Inches(4.5), Inches(2.5), Inches(4.3), Inches(3.0), NAVY)
add_text(s, Inches(4.5), Inches(2.65), Inches(4.3), Inches(0.5),
         "Your Company Folder", font_size=18, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

folder_items = [
    "company-context/     ICP, positioning",
    "skills/              outreach templates",
    "workflows/           enrichment steps",
    "memory/              lessons learned",
    "CLAUDE.md            operating instructions",
]
add_multiline(s, Inches(4.9), Inches(3.2), Inches(3.6), Inches(2.2),
              folder_items, font_size=11, font_name=FONT_MONO, color=RGBColor(0xAA, 0xAA, 0xBB),
              line_spacing=1.8)

# Agent label on top
agent_shape = add_rect(s, Inches(5.3), Inches(2.1), Inches(2.6), Inches(0.45), GOLD)
add_text(s, Inches(5.3), Inches(2.12), Inches(2.6), Inches(0.4),
         "AI Agent", font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Left side: Connected tools (CRM)
left_tools = [("HubSpot", 2.8), ("Salesforce", 3.5), ("Google Drive", 4.2)]
for name, y in left_tools:
    add_rect(s, Inches(0.8), Inches(y), Inches(2.8), Inches(0.5), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(0.8), Inches(y + 0.03), Inches(2.8), Inches(0.45),
             name, font_size=13, color=NAVY, alignment=PP_ALIGN.CENTER)

# Arrows left -> center
for y in [3.0, 3.7, 4.4]:
    add_text(s, Inches(3.7), Inches(y - 0.05), Inches(0.8), Inches(0.4),
             "\u2194", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# Right side: Web tools
right_tools = [("Web Search", 2.8), ("Google Maps", 3.5), ("LinkedIn", 4.2)]
for name, y in right_tools:
    add_rect(s, Inches(9.7), Inches(y), Inches(2.8), Inches(0.5), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(9.7), Inches(y + 0.03), Inches(2.8), Inches(0.45),
             name, font_size=13, color=NAVY, alignment=PP_ALIGN.CENTER)

# Arrows center -> right
for y in [3.0, 3.7, 4.4]:
    add_text(s, Inches(8.8), Inches(y - 0.05), Inches(0.8), Inches(0.4),
             "\u2194", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# Bottom callout
add_text(s, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.9),
         "All-in-one. No tab-switching. No copy-pasting.\nThe agent reads your context, executes across all connected tools, and delivers structured output.",
         font_size=14, color=NAVY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)


# ════════════════════════════════════════════════════════
# SLIDE 6 — Connects to Your Stack
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Connects to your stack",
         font_size=36, font_name=FONT_HEADING, color=WHITE, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "No migration. No new software. ZPT plugs into what you already use.",
         font_size=16, color=MUTED)

# Tool logo placeholders in a grid
tools = ["HubSpot", "Salesforce", "Google Drive", "Notion", "Slack",
         "LinkedIn", "Google Maps", "ZeroBounce"]

cols = 4
for i, tool in enumerate(tools):
    col = i % cols
    row = i // cols
    x = Inches(1.8 + col * 2.6)
    y = Inches(2.5 + row * 2.0)

    add_rect(s, x, y, Inches(2.0), Inches(1.5), NAVY_LIGHT, border_color=RGBColor(0x44, 0x44, 0x66))
    add_text(s, x, y + Inches(0.5), Inches(2.0), Inches(0.5),
             tool, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER, bold=True)

# NOTE: This slide uses placeholder boxes. Wybe should add actual logos.
add_text(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
         "[ Replace boxes with actual tool logos from /public/logos/ ]",
         font_size=11, color=RGBColor(0x66, 0x66, 0x88), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 7 — How We Get Started
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "How we get started",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

steps = [
    ("1", "Collect company context",
     "We learn your ICP, positioning, current process, and examples of good-fit accounts."),
    ("2", "Connect to your CRM",
     "We hook into HubSpot, Salesforce, or whatever your team already uses."),
    ("3", "Start delivering results",
     "ZPT begins sourcing, enriching, and qualifying prospects. Output flows back through your CRM."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)

    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    add_text(s, x + Inches(1.3), y + Inches(0.1), Inches(0.8), Inches(0.6),
             num, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text(s, x, y + Inches(1.1), Inches(3.6), Inches(0.5),
             title, font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text(s, x + Inches(0.1), y + Inches(1.7), Inches(3.4), Inches(1.5),
             desc, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Arrow between steps
    if i < 2:
        add_text(s, x + Inches(3.6), y + Inches(0.15), Inches(0.5), Inches(0.6),
                 "\u2192", font_size=28, color=GOLD, alignment=PP_ALIGN.CENTER)

# Callout box
add_rect(s, Inches(3.2), Inches(5.0), Inches(6.8), Inches(0.8), CREAM)
add_text(s, Inches(3.2), Inches(5.05), Inches(6.8), Inches(0.7),
         "You get the output. We run the system.",
         font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Trial note
add_text(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
         "Start with 100 contacts enriched for free.",
         font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 8 — Already Operational
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Already operational",
         font_size=36, font_name=FONT_HEADING, color=WHITE, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

cases = [
    ("A logistics company",
     "20,000 contacts enriched, verified, and cleaned.",
     "The sales team now reaches significantly more companies with higher conversion rates."),
    ("Blankert Shortlease",
     "Structured prospect sourcing in the Dutch automotive leasing market.",
     "Ongoing enrichment and outreach support across directories, industry events, and LinkedIn."),
    ("Financial Analytics",
     "1,300+ contacts prospected across institutional investors.",
     "Full profiles with decision-maker identification and company intelligence."),
]

for i, (company, stat, detail) in enumerate(cases):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)

    add_rect(s, x, y, Inches(3.6), Inches(4.2), NAVY_LIGHT, border_color=RGBColor(0x44, 0x44, 0x66))

    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.6),
             company, font_size=22, color=GOLD, bold=True, font_name=FONT_HEADING)

    add_text(s, x + Inches(0.3), y + Inches(1.1), Inches(3.0), Inches(1.2),
             stat, font_size=16, color=WHITE, bold=True, line_spacing=1.4)

    add_text(s, x + Inches(0.3), y + Inches(2.5), Inches(3.0), Inches(1.2),
             detail, font_size=13, color=MUTED, line_spacing=1.5)


# ════════════════════════════════════════════════════════
# SLIDE 9 — Who This Is For
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Who this is for",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

profiles = [
    ("Inconsistent enrichment",
     "Companies with sales data but no reliable, repeatable way to enrich and verify it."),
    ("Spreadsheet chaos",
     "Teams working from spreadsheets or a CRM that needs serious cleanup before it's useful."),
    ("More pipeline, no new hires",
     "Firms that want to grow their pipeline without the cost and time of hiring more people."),
    ("Lean teams, big goals",
     "Small sales teams that need enterprise-level output from their existing process."),
]

for i, (title, desc) in enumerate(profiles):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 2.3)

    add_card(s, x, y, Inches(5.6), Inches(1.9), title, desc,
             fill=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD),
             title_size=18, body_size=14)


# ════════════════════════════════════════════════════════
# SLIDE 10 — CTA
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
         "Ready to see what ZPT\ncan do for your sales?",
         font_size=40, font_name=FONT_HEADING, color=WHITE,
         alignment=PP_ALIGN.CENTER, line_spacing=1.3)

add_gold_divider(s, Inches(5.5), Inches(3.2), Inches(2.3))

# CTA button
add_rect(s, Inches(4.8), Inches(3.7), Inches(3.7), Inches(0.8), GOLD)
add_text(s, Inches(4.8), Inches(3.78), Inches(3.7), Inches(0.65),
         "Book a Call", font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
         "request@zpteam.ai",
         font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
         "Start with 100 contacts enriched for free.",
         font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/Users/wybeharms/Sites/zpt/dev/public/presentation/ZPT_Sales_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
