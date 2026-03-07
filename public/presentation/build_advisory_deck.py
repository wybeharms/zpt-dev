"""
Build ZPT Advisory Deck — 10 slides, ZPT brand (Navy/Gold/Off-white).
Run: python3 build_advisory_deck.py
Output: ZPT_Advisory_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Brand tokens --
NAVY = RGBColor(0x0C, 0x0C, 0x28)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x88, 0x88, 0x99)
LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0xA8)
NAVY_LIGHT = RGBColor(0x1A, 0x1A, 0x3E)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)

FONT_HEADING = "Cormorant Garamond"
FONT_BODY = "Plus Jakarta Sans"
FONT_MONO = "Source Code Pro"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # blank layout


# -- Helpers --

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             font_name=FONT_BODY, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  font_name=FONT_BODY, color=WHITE, alignment=PP_ALIGN.LEFT,
                  line_spacing=1.4):
    """lines: list of (text, bold) tuples or strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, bold = line
        else:
            text, bold = line, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_card(slide, left, top, width, height, title, body, fill=WHITE,
             title_color=NAVY, body_color=MUTED, border_color=None,
             title_size=16, body_size=13):
    """Add a rounded card with title + body text."""
    card = add_rect(slide, left, top, width, height, fill, border_color)
    add_text(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.5),
             title, font_size=title_size, font_name=FONT_BODY, color=title_color, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
             body, font_size=body_size, font_name=FONT_BODY, color=body_color)
    return card


def add_gold_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


# ================================================================
# SLIDE 1 -- Title
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "Your company knows AI is powerful.",
         font_size=44, font_name=FONT_HEADING, color=WHITE, bold=False,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(2.9), Inches(10), Inches(1),
         "We make it happen.",
         font_size=44, font_name=FONT_HEADING, color=GOLD, bold=False,
         alignment=PP_ALIGN.CENTER)

add_gold_divider(s, Inches(5.5), Inches(4.2), Inches(2.3))

add_text(s, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
         "zpteam.ai  |  request@zpteam.ai",
         font_size=14, font_name=FONT_BODY, color=MUTED,
         alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2 -- The Overwhelm
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "The overwhelm",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

# Left: buzzword wall
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
add_text(s, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
         "What your team sees every week",
         font_size=18, font_name=FONT_BODY, color=NAVY, bold=True)

buzzwords = [
    "MCPs", "skills", "sub-agents", "plugins", "CLIs",
    "memory systems", "agent frameworks", "prompt engineering",
    "RAG", "fine-tuning", "embeddings", "vector databases",
    "function calling", "tool use", "agentic workflows",
    "context windows", "system prompts", "guardrails",
]

# Scatter buzzwords in a word-cloud style grid
row_y = Inches(2.5)
col_x = Inches(1.3)
for i, word in enumerate(buzzwords):
    add_text(s, col_x, row_y, Inches(2.2), Inches(0.4),
             word, font_size=13, font_name=FONT_MONO,
             color=RGBColor(0x99, 0x99, 0xAA))
    col_x += Inches(1.7)
    if col_x > Inches(4.8):
        col_x = Inches(1.3)
        row_y += Inches(0.45)

add_text(s, Inches(1.1), Inches(5.8), Inches(5), Inches(0.6),
         "[ Wybe: style this as a visual buzzword wall in PowerPoint ]",
         font_size=10, color=RGBColor(0xBB, 0xBB, 0xCC))

# Right: ZPT's answer
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), NAVY)
add_text(s, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5),
         "What ZPT does",
         font_size=18, font_name=FONT_BODY, color=GOLD, bold=True)

add_text(s, Inches(7.3), Inches(2.5), Inches(4.8), Inches(2.0),
         "We boil it down to what truly matters:\n\nOne folder. One agent.\nConnected to your tools.",
         font_size=18, font_name=FONT_BODY, color=WHITE, line_spacing=1.6)

add_text(s, Inches(7.3), Inches(4.5), Inches(4.8), Inches(1.5),
         "Your team doesn't need to understand the technology.\nThey just need a setup that works.",
         font_size=14, font_name=FONT_BODY, color=MUTED, line_spacing=1.5)

add_rect(s, Inches(7.3), Inches(5.6), Inches(4.8), Inches(0.7), NAVY_LIGHT)
add_text(s, Inches(7.6), Inches(5.68), Inches(4.3), Inches(0.5),
         "First principles, not feature chasing.",
         font_size=15, font_name=FONT_BODY, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 3 -- From Chatbot to Agent
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "From chatbot to agent",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "The same AI your team already uses, but now it actually does work.",
         font_size=16, color=MUTED)

# Left card: Chatbot
add_rect(s, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.2), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
add_text(s, Inches(1.1), Inches(2.4), Inches(5), Inches(0.5),
         "Chatbot", font_size=22, color=NAVY, bold=True)
add_text(s, Inches(1.1), Inches(3.0), Inches(5), Inches(0.5),
         "Text in, text out.", font_size=16, color=MUTED, bold=True)

chatbot_lines = [
    "Answer questions",
    "Summarize text",
    "Write drafts",
    "Brainstorm ideas",
]
add_multiline(s, Inches(1.1), Inches(3.6), Inches(5), Inches(2.5),
              [f"+  {l}" for l in chatbot_lines],
              font_size=14, color=MUTED, line_spacing=1.8)

add_rect(s, Inches(1.1), Inches(5.4), Inches(4.8), Inches(0.7), CREAM)
add_text(s, Inches(1.1), Inches(5.45), Inches(4.8), Inches(0.6),
         "It talks.", font_size=16, color=NAVY, bold=True,
         alignment=PP_ALIGN.CENTER)

# Right card: Agent
add_rect(s, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.2), NAVY)
add_text(s, Inches(7.3), Inches(2.4), Inches(5), Inches(0.5),
         "Agent", font_size=22, color=GOLD, bold=True)
add_text(s, Inches(7.3), Inches(3.0), Inches(5), Inches(0.5),
         "Instruction in, actions out.", font_size=16, color=WHITE, bold=True)

agent_lines = [
    "Read files and documents",
    "Search the web",
    "Update your CRM",
    "Chain multiple steps together",
]
add_multiline(s, Inches(7.3), Inches(3.6), Inches(5), Inches(2.5),
              [f"+  {l}" for l in agent_lines],
              font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.8)

add_rect(s, Inches(7.3), Inches(5.4), Inches(4.8), Inches(0.7), GOLD)
add_text(s, Inches(7.3), Inches(5.45), Inches(4.8), Inches(0.6),
         "It works.", font_size=16, color=NAVY, bold=True,
         alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 4 -- How It Works / The Folder
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "How it works",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "One local folder. The AI agent reads it, connects to your tools, and executes tasks.",
         font_size=16, color=MUTED)

# Center: The folder
add_rect(s, Inches(4.2), Inches(2.4), Inches(5.0), Inches(3.2), NAVY)
add_text(s, Inches(4.2), Inches(2.55), Inches(5.0), Inches(0.5),
         "Your Company Folder", font_size=18, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

folder_items = [
    "company-context/     ICP, positioning, overview",
    "skills/              templates, quality standards",
    "workflows/           step-by-step procedures",
    "memory/              lessons learned over time",
    "mcp/                 tool configurations",
    "CLAUDE.md            operating instructions",
]
add_multiline(s, Inches(4.6), Inches(3.1), Inches(4.3), Inches(2.4),
              folder_items, font_size=11, font_name=FONT_MONO, color=RGBColor(0xAA, 0xAA, 0xBB),
              line_spacing=1.7)

# Agent label on top
add_rect(s, Inches(5.5), Inches(2.0), Inches(2.6), Inches(0.45), GOLD)
add_text(s, Inches(5.5), Inches(2.02), Inches(2.6), Inches(0.4),
         "AI Agent", font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Left side: Connected tools
left_tools = [("HubSpot", 2.7), ("Google Drive", 3.4), ("Slack", 4.1)]
for name, y in left_tools:
    add_rect(s, Inches(0.8), Inches(y), Inches(2.6), Inches(0.5), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(0.8), Inches(y + 0.03), Inches(2.6), Inches(0.45),
             name, font_size=13, color=NAVY, alignment=PP_ALIGN.CENTER)

# Arrows left -> center
for y in [2.9, 3.6, 4.3]:
    add_text(s, Inches(3.5), Inches(y - 0.05), Inches(0.7), Inches(0.4),
             "\u2194", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# Right side: More tools
right_tools = [("Excel", 2.7), ("PowerPoint", 3.4), ("Notion", 4.1)]
for name, y in right_tools:
    add_rect(s, Inches(9.9), Inches(y), Inches(2.6), Inches(0.5), WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(9.9), Inches(y + 0.03), Inches(2.6), Inches(0.45),
             name, font_size=13, color=NAVY, alignment=PP_ALIGN.CENTER)

# Arrows center -> right
for y in [2.9, 3.6, 4.3]:
    add_text(s, Inches(9.2), Inches(y - 0.05), Inches(0.7), Inches(0.4),
             "\u2194", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# Bottom callout
add_rect(s, Inches(2.5), Inches(6.0), Inches(8.3), Inches(1.0), CREAM)
add_text(s, Inches(2.5), Inches(6.05), Inches(8.3), Inches(0.45),
         "Same concept as OpenClaw (250K+ GitHub stars), but fully owned by you and much simpler.",
         font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(6.45), Inches(8.3), Inches(0.45),
         "The AI provider is interchangeable. Claude, ChatGPT, open-source \u2014 your folder stays the same.",
         font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 5 -- What We Offer (4 service cards)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "What we offer",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Every company is different. Choose what you need. You own everything we build.",
         font_size=16, color=MUTED)

services = [
    ("Educate your team",
     "Presentations, demos, Q&A tailored to the company's industry and workflows."),
    ("Set up accounts and tools",
     "Configure Claude, ChatGPT, desktop apps, and get your team up and running."),
    ("Build the company setup",
     "Structured folder with company context, workflows, templates, and quality standards."),
    ("Integrate your tools",
     "Wire up MCPs and integrations to existing tools like CRM, drives, and messaging."),
]

for i, (title, desc) in enumerate(services):
    x = Inches(0.8 + i * 3.1)
    add_card(s, x, Inches(2.2), Inches(2.8), Inches(3.5), title, desc,
             fill=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD),
             title_size=16, body_size=13)

# Icon placeholders
icons = ["Presentation", "Claude", "Folder", "Plug"]
for i, icon in enumerate(icons):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x + Inches(0.3), Inches(4.7), Inches(0.6), Inches(0.35), GOLD)
    add_text(s, x + Inches(0.3), Inches(4.72), Inches(0.6), Inches(0.3),
             icon[0], font_size=11, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 6 -- What Your Team Gets (6 use cases)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "What your team gets",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

use_cases = [
    ("Build presentations from scratch",
     "Have Claude create an entire Google Slides deck based on your company materials, style, and previous client work."),
    ("Pull from multiple tools",
     "Find the right meeting notes, summarize them, and add the outcome to Notion or another system your team already uses."),
    ("Standardize repetitive workflows",
     "Turn recurring manual work into a repeatable workflow your team can actually use without re-explaining everything each time."),
    ("Create better outputs",
     "Generate summaries, drafts, reports, and updates in a format that matches how your company already works."),
    ("Connect context across tools",
     "Use information from notes, documents, spreadsheets, and internal knowledge in one place instead of scattered chat histories."),
    ("Onboard team members faster",
     "Give new hires access to the same company knowledge, templates, and workflows from day one."),
]

for i, (title, desc) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.6)
    add_card(s, x, y, Inches(3.7), Inches(2.2), title, desc,
             fill=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD),
             title_size=15, body_size=12)

# Bottom callout
add_rect(s, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.5), CREAM)
add_text(s, Inches(2.5), Inches(6.92), Inches(8.3), Inches(0.45),
         "Or: we identify your most repetitive workflows and automate them.",
         font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 7 -- How a ZPT Engagement Works
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, OFF_WHITE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "How a ZPT Advisory engagement works",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "We visit your company, understand how the team works, and build the right AI setup.",
         font_size=16, color=MUTED)

# Placeholder for helicopter/SWAT visual
add_text(s, Inches(0.8), Inches(2.1), Inches(4), Inches(0.4),
         "[ Wybe: add helicopter/SWAT team visual from the website here ]",
         font_size=10, color=RGBColor(0xBB, 0xBB, 0xCC))

steps = [
    ("1", "Understand workflows",
     "Learn the team's tools, processes, and pain points."),
    ("2", "Educate the team",
     "Presentations and hands-on demos tailored to the company's workflows."),
    ("3", "Set up accounts & tools",
     "Configure AI tools and desktop apps for every team member."),
    ("4", "Build the company folder",
     "Optionally create a tailored setup with context, skills, and integrations."),
    ("5", "Train the team",
     "Hands-on training so the team works with AI as a co-worker."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.55)
    y = Inches(2.5)

    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    add_text(s, x + Inches(0.85), y + Inches(0.08), Inches(0.7), Inches(0.55),
             num, font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text(s, x, y + Inches(1.0), Inches(2.4), Inches(0.5),
             title, font_size=15, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text(s, x + Inches(0.05), y + Inches(1.5), Inches(2.3), Inches(1.2),
             desc, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Arrow between steps
    if i < 4:
        add_text(s, x + Inches(2.35), y + Inches(0.12), Inches(0.3), Inches(0.5),
                 "\u2192", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# Callout box
add_rect(s, Inches(3.2), Inches(5.5), Inches(6.8), Inches(0.7), CREAM)
add_text(s, Inches(3.2), Inches(5.55), Inches(6.8), Inches(0.6),
         "Everything ZPT builds stays with the company. No lock-in, no dependencies.",
         font_size=15, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 8 -- What You Own
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "What you own",
         font_size=36, font_name=FONT_HEADING, color=NAVY, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

own_cards = [
    ("100% company-owned",
     "The entire system runs on a local folder on your machine. No data leaves your network. No vendor lock-in. The most private AI setup possible."),
    ("Open-source compatible",
     "Concerned about sharing data with AI labs? Run the system on open-source models (Llama, Mistral, etc.). Same folder, same workflow, full control."),
    ("Built on Claude Code architecture",
     "The same local folder system used by developers worldwide (CLAUDE.md, skills/, memory/), but purpose-built for business operations instead of code."),
    ("Plugs into what you already use",
     "Claude in your spreadsheets. Claude in your presentations. Connect to Excel, PowerPoint, Google Drive, Slack, HubSpot, and more. No migration required."),
]

for i, (title, desc) in enumerate(own_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.6 + row * 2.5)
    add_card(s, x, y, Inches(5.6), Inches(2.1), title, desc,
             fill=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD),
             title_size=17, body_size=13)

# Interchangeable provider note
add_text(s, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
         "The AI provider is interchangeable: Claude, ChatGPT, or open-source. Your company folder stays the same.",
         font_size=14, color=NAVY, alignment=PP_ALIGN.CENTER, bold=True)


# ================================================================
# SLIDE 9 -- Advisory Tiers
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Advisory tiers",
         font_size=36, font_name=FONT_HEADING, color=WHITE, bold=False)
add_gold_divider(s, Inches(0.8), Inches(1.2), Inches(2))

add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Choose what fits \u2014 pricing on request.",
         font_size=16, color=MUTED)

tiers = [
    ("Educate", "Half-day", "Your team understands AI",
     ["What AI agents are and how they work",
      "Industry-specific use cases and demos",
      "Hands-on Q&A with your team's workflows"]),
    ("Equip", "1-3 days", "Your team uses AI daily",
     ["Everything in Educate, plus:",
      "Set up accounts and desktop apps",
      "Configure company context and templates",
      "Connect existing tools via integrations"]),
    ("Embed", "Ongoing", "AI becomes part of how the company works",
     ["Everything in Equip, plus:",
      "Build and maintain company folder",
      "Identify and automate key workflows",
      "Ongoing training and support"]),
]

for i, (name, time, tagline, features) in enumerate(tiers):
    x = Inches(0.8 + i * 4.1)
    is_highlight = (i == 1)
    border = GOLD if is_highlight else RGBColor(0x44, 0x44, 0x66)

    add_rect(s, x, Inches(2.0), Inches(3.6), Inches(4.8), NAVY_LIGHT, border_color=border)
    if is_highlight:
        # Extra gold border thickness
        pass

    # Tier name
    add_text(s, x, Inches(2.2), Inches(3.6), Inches(0.6),
             name, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Time badge
    add_rect(s, x + Inches(1.0), Inches(2.85), Inches(1.6), Inches(0.4), GOLD)
    add_text(s, x + Inches(1.0), Inches(2.87), Inches(1.6), Inches(0.35),
             time, font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Tagline
    add_text(s, x + Inches(0.3), Inches(3.45), Inches(3.0), Inches(0.5),
             tagline, font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER, bold=True)

    # Features
    add_multiline(s, x + Inches(0.3), Inches(4.1), Inches(3.0), Inches(2.5),
                  [f"+  {f}" for f in features],
                  font_size=12, color=RGBColor(0xAA, 0xAA, 0xBB), line_spacing=1.7)


# ================================================================
# SLIDE 10 -- CTA
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
         "Ready to see what AI\ncan do for your company?",
         font_size=40, font_name=FONT_HEADING, color=WHITE,
         alignment=PP_ALIGN.CENTER, line_spacing=1.3)

add_gold_divider(s, Inches(5.5), Inches(3.2), Inches(2.3))

# CTA button
add_rect(s, Inches(4.8), Inches(3.7), Inches(3.7), Inches(0.8), GOLD)
add_text(s, Inches(4.8), Inches(3.78), Inches(3.7), Inches(0.65),
         "Book a Call", font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
         "request@zpteam.ai",
         font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
         "zpteam.ai/advisory",
         font_size=14, color=RGBColor(0x66, 0x66, 0x88), alignment=PP_ALIGN.CENTER)


# -- Save --
output_path = "/Users/wybeharms/Sites/zpt/dev/public/presentation/ZPT_Advisory_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
